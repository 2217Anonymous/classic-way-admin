"""Generate Valaiyagam client PPT and project tracking Excel."""

from datetime import date
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from tracker_stories import (
    ALL_STORIES,
    E2E_WORKFLOWS,
    START,
    TIMELINE_PHASES,
    week_date,
)

ROOT = Path(__file__).resolve().parent


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
    shape.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0x14, 0xB8, 0xA6)
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    sub = slide.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11.5), Inches(1.2))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(18)
    sp.font.color.rgb = RGBColor(0x47, 0x55, 0x69)


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], footer: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
    bg.line.fill.background()

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.55), Inches(0.12), Inches(0.55)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x14, 0xB8, 0xA6)
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.5), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    body = slide.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(11.3), Inches(5.3))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
        p.space_after = Pt(10)

    if footer:
        fb = slide.shapes.add_textbox(Inches(1.0), Inches(6.9), Inches(11.3), Inches(0.35))
        fp = fb.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(12)
        fp.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Valaiyagam E-Commerce Platform",
        "Architecture · UI Standards · Payments · Courier · Tracking · Timeline\nClient Presentation · July 2026",
    )

    add_bullets_slide(
        prs,
        "Agenda",
        [
            "1. Product vision and business outcomes",
            "2. End-to-end architecture (Admin + Shopping + PostgreSQL)",
            "3. UI / UX standards (light glass theme)",
            "4. Payment integration workflow",
            "5. Courier integration workflow",
            "6. Order tracking system",
            "7. Delivery timeline and milestones",
            "8. Git / story tracking process",
            "9. Next steps and client inputs needed",
        ],
    )

    add_bullets_slide(
        prs,
        "Product Vision",
        [
            "One platform for storefront shopping and admin operations",
            "Secure role-based admin for catalog, orders, users, and fulfillment",
            "Trusted checkout with payment gateway + COD options",
            "Reliable courier booking, labels, and live shipment tracking",
            "Mobile-first customer and admin experiences",
            "Dockerized, migration-ready foundation for long-term growth",
        ],
        "Foundation already live: Auth, Users, Roles, Docker, Alembic, Glass Admin UI",
    )

    add_bullets_slide(
        prs,
        "Architecture Overview",
        [
            "Frontend: Next.js + React + Tailwind + Redux Toolkit (Admin + Shopping theme)",
            "Backend: FastAPI layered design (Route → Service → Repository → Model)",
            "Database: PostgreSQL 16 with Alembic versioned migrations",
            "Infrastructure: Docker Compose (admin/shopping frontends + backends + postgres)",
            "Integrations: Payment gateway APIs + Courier partner adapters",
            "Security: JWT auth (staff + customer), Argon2 passwords, RBAC, webhook signatures",
        ],
    )

    add_bullets_slide(
        prs,
        "UI Standards — Visual System",
        [
            "Light theme only — soft #f5f8fc backgrounds (no dark shells)",
            "Glassmorphism panels: translucent white, blur, soft borders",
            "Primary accent: teal → cyan gradient buttons and active states",
            "Typography: clear hierarchy, generous spacing, readable contrast",
            "Icons: simple line icons for navigation and actions",
            "Motion: subtle hover lift / shadow — never noisy animation",
        ],
    )

    add_bullets_slide(
        prs,
        "UI Standards — Layout & Interaction",
        [
            "Desktop: glass sidebar + content workspace",
            "Mobile: sticky bottom navigation + card lists (not dense tables)",
            "Create / Edit always open in glass modal forms",
            "Delete always requires a confirmation modal (never browser alerts)",
            "One primary job per screen/section",
            "Inline error banners from API; empty states with clear CTA",
        ],
    )

    add_bullets_slide(
        prs,
        "UI Standards — Responsive Rules",
        [
            "Breakpoints cover phone, tablet, and desktop admin use",
            "Touch targets sized for thumbs on mobile actions",
            "Modals become bottom sheets on small screens",
            "Tables transform into stacked cards under tablet width",
            "Esc key closes dialogs; focus remains usable",
            "Consistent glass components reused across all modules",
        ],
    )

    add_bullets_slide(
        prs,
        "Admin ↔ Shopping Sync",
        [
            "Admin publish product → appears on shopping site (shared PostgreSQL)",
            "Shopping place order → visible instantly in Admin Orders",
            "Admin mark paid / ship → customer Track Order updates",
            "Architecture: admin :3000/:8000 · shopping :3001/:8001 · same DB",
            "Theme: BlueBerry fashion template wired to live shopping APIs",
            "Stories: VL-* foundation + VS-* shopping sync in PROJECT_TRACKING.xlsx",
        ],
    )

    add_bullets_slide(
        prs,
        "Admin Modules Roadmap",
        [
            "Done: Authentication, Users, Roles, Catalog, Inventory, Orders",
            "In progress: Payments gateway hardening, Courier AWB, Tracking polish",
            "Shopping sync epics: Theme, Catalog reflect, Commerce, Admin reflect, Fulfill",
            "Then: Customer console, review moderation, tax in checkout, CMS banners",
            "Each module follows the same UI + API layering standards",
        ],
    )

    add_bullets_slide(
        prs,
        "Payment Integration",
        [
            "Customer checkout creates order in PENDING_PAYMENT state",
            "Backend creates gateway order/intent (Razorpay primary; Stripe optional)",
            "Customer pays via hosted checkout / SDK — no card storage in our DB",
            "Gateway webhook verifies signature → marks order PAID",
            "Failure/timeout → PAYMENT_FAILED with safe retry",
            "Admin can issue refunds; all events audited in payment_events",
            "COD supported with courier confirmation before settlement",
        ],
    )

    add_bullets_slide(
        prs,
        "Payment — Status Lifecycle",
        [
            "PENDING_PAYMENT → PAID → (eligible for fulfillment)",
            "PENDING_PAYMENT → PAYMENT_FAILED → retry / cancel",
            "PAID → REFUND_REQUESTED → PARTIALLY_REFUNDED / REFUNDED",
            "Webhook idempotency prevents double capture",
            "Nightly reconciliation job catches missed webhooks",
            "Settlement reports available in admin finance views",
        ],
    )

    add_bullets_slide(
        prs,
        "Courier Integration",
        [
            "Adapter pattern supports multiple partners under one service",
            "Candidates: Delhivery, Shiprocket, BlueDart + manual fallback",
            "After PAID (or approved COD): create shipment → get AWB",
            "Generate shipping label and schedule pickup",
            "Serviceability check by pincode/weight before booking",
            "Cancel / reassign shipment with audit trail",
            "Credentials stored as secure courier_accounts configuration",
        ],
    )

    add_bullets_slide(
        prs,
        "Courier — Operational Flow",
        [
            "1. Validate address, package weight, and serviceability",
            "2. Apply routing rules (cost, SLA, zone)",
            "3. Call partner create_shipment API",
            "4. Persist AWB + label URL on shipment record",
            "5. Print label / handoff to warehouse packing",
            "6. Sync status via partner webhook or scheduled poll",
            "7. Handle exceptions: RTO, address issue, lost package",
        ],
    )

    add_bullets_slide(
        prs,
        "Tracking System",
        [
            "Unified timeline for customer and admin",
            "Stages: Placed → Paid → Packed → Shipped → Out for delivery → Delivered",
            "Customer track page: Order ID + verified phone/email (or login)",
            "Admin order detail shows payment events + shipment checkpoints",
            "Notifications on key milestones (email/SMS/WhatsApp)",
            "Exception queue for delayed / failed deliveries",
            "Manual override allowed for ops with mandatory reason log",
        ],
    )

    add_bullets_slide(
        prs,
        "End-to-End Order Journey",
        [
            "Browse catalog → Add to cart → Checkout address",
            "Choose payment method → Pay / COD confirm",
            "Warehouse packs order → Courier booked → AWB issued",
            "Customer receives tracking link + live timeline updates",
            "Delivery / RTO / refund handled with status sync",
            "Reports update sales, fulfillment SLA, and settlements",
        ],
    )

    add_bullets_slide(
        prs,
        "16-Week Delivery Timeline",
        [
            "W1–W2  P0 Foundation — Auth, Users, Roles, Docker (DONE)",
            "W3–W5  P1 Catalog — Products, categories, inventory basics",
            "W6–W8  P2 Orders — Cart, checkout, order lifecycle",
            "W9–W10 P3 Payments — Gateway, webhooks, refunds, COD",
            "W11–W12 P4 Courier — Adapters, AWB, labels, pickup",
            "W13–W14 P5 Tracking & Notify — Timeline + alerts",
            "W15–W16 P6 Harden & Launch — UAT, performance, go-live",
        ],
    )

    add_bullets_slide(
        prs,
        "Milestones for Client Sign-off",
        [
            "M1 Foundation demo — complete",
            "M2 Catalog demo — publish products",
            "M3 Checkout demo — place unpaid order",
            "M4 Payment demo — paid order path",
            "M5 Courier demo — AWB + label from paid order",
            "M6 Tracking demo — customer timeline live",
            "M7 Go-live readiness checklist approved",
        ],
    )

    add_bullets_slide(
        prs,
        "Git & Story Tracking Process",
        [
            "Branches: main · develop · feat/* · fix/* · hotfix/*",
            "Each story maps to a feature branch and PR",
            "Excel tracker stores: ID, Date, Phase, Story, Branch, Status, Owner",
            "Definition of Done: code + review + migration (if any) + UAT note",
            "Weekly status pulled directly from PROJECT_TRACKING.xlsx",
            "No direct commits to main — PR + review required",
        ],
    )

    add_bullets_slide(
        prs,
        "Client Inputs Needed",
        [
            "Brand logo, colors confirmation, and product sample data",
            "Payment merchant account (Razorpay/Stripe) keys for staging",
            "Courier partner accounts and warehouse pickup addresses",
            "Tax rules, shipping zones, and COD eligibility policy",
            "Notification channels (email domain / SMS provider)",
            "UAT users and preferred weekly review slot",
        ],
    )

    add_bullets_slide(
        prs,
        "Summary & Ask",
        [
            "Solid technical foundation is already running in Docker",
            "UI standard is fixed: light glass, modal edit, confirm deletes, mobile-first",
            "Payment, courier, and tracking are designed as first-class modules",
            "16-week plan delivers a launchable MVP with clear milestones",
            "Ask: approve timeline, confirm gateway + courier partners, start P1 Catalog",
        ],
        "Thank you — Valaiyagam Commerce Team",
    )

    prs.save(path)


def style_header(ws, row: int, cols: int) -> None:
    fill = PatternFill("solid", fgColor="0F766E")
    font = Font(color="FFFFFF", bold=True, size=11)
    thin = Border(
        left=Side(style="thin", color="CBD5E1"),
        right=Side(style="thin", color="CBD5E1"),
        top=Side(style="thin", color="CBD5E1"),
        bottom=Side(style="thin", color="CBD5E1"),
    )
    for col in range(1, cols + 1):
        cell = ws.cell(row=row, column=col)
        cell.fill = fill
        cell.font = font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = thin


def autosize(ws, widths: dict[int, int]) -> None:
    for col, width in widths.items():
        ws.column_dimensions[get_column_letter(col)].width = width


def status_fill(status: str) -> PatternFill | None:
    mapping = {
        "Done": "D1FAE5",
        "In Progress": "FEF3C7",
        "Planned": "E0F2FE",
        "Blocked": "FEE2E2",
    }
    color = mapping.get(status)
    return PatternFill("solid", fgColor=color) if color else None


def build_excel(path: Path) -> None:
    wb = Workbook()

    # --- Stories sheet ---
    stories = wb.active
    stories.title = "Stories"

    headers = [
        "ID",
        "Date",
        "Phase",
        "Week",
        "Module",
        "Story Title",
        "User Story",
        "Git Feature / Branch",
        "Branch Type",
        "Priority",
        "Status",
        "Owner",
        "Estimate (Days)",
        "Depends On",
        "Acceptance Criteria",
        "Notes",
    ]
    stories.append(headers)
    style_header(stories, 1, len(headers))

    rows = list(ALL_STORIES)

    thin = Border(
        left=Side(style="thin", color="E2E8F0"),
        right=Side(style="thin", color="E2E8F0"),
        top=Side(style="thin", color="E2E8F0"),
        bottom=Side(style="thin", color="E2E8F0"),
    )

    for row in rows:
        stories.append(list(row))
        r = stories.max_row
        for c in range(1, len(headers) + 1):
            cell = stories.cell(row=r, column=c)
            cell.border = thin
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            if c == 2:
                cell.number_format = "YYYY-MM-DD"
        status_cell = stories.cell(row=r, column=11)
        fill = status_fill(str(status_cell.value))
        if fill:
            status_cell.fill = fill

    autosize(
        stories,
        {
            1: 10,
            2: 12,
            3: 16,
            4: 8,
            5: 14,
            6: 40,
            7: 55,
            8: 38,
            9: 12,
            10: 10,
            11: 12,
            12: 12,
            13: 14,
            14: 14,
            15: 48,
            16: 28,
        },
    )
    stories.freeze_panes = "A2"
    stories.auto_filter.ref = f"A1:P{stories.max_row}"
    stories.row_dimensions[1].height = 30

    # --- E2E Workflows sheet ---
    workflows = wb.create_sheet("E2E Workflows")
    wf_headers = [
        "Workflow ID",
        "Name",
        "Starts In",
        "Ends In",
        "Admin / Source Steps",
        "Shopping / Target Steps",
        "Linked Stories",
        "Priority",
        "Status",
    ]
    workflows.append(wf_headers)
    style_header(workflows, 1, len(wf_headers))
    for item in E2E_WORKFLOWS:
        workflows.append(list(item))
        r = workflows.max_row
        for c in range(1, len(wf_headers) + 1):
            workflows.cell(row=r, column=c).alignment = Alignment(
                vertical="top", wrap_text=True
            )
            workflows.cell(row=r, column=c).border = thin
        fill = status_fill(str(item[-1]))
        if fill:
            workflows.cell(row=r, column=9).fill = fill
    autosize(
        workflows,
        {1: 12, 2: 18, 3: 12, 4: 12, 5: 55, 6: 55, 7: 36, 8: 10, 9: 12},
    )
    workflows.freeze_panes = "A2"
    workflows.row_dimensions[1].height = 30
    for r in range(2, workflows.max_row + 1):
        workflows.row_dimensions[r].height = 60

    # --- Timeline sheet ---
    timeline = wb.create_sheet("Timeline")
    timeline_headers = [
        "Phase",
        "Week Start",
        "Week End",
        "Focus",
        "Milestone",
        "Exit Criteria",
        "Status",
    ]
    timeline.append(timeline_headers)
    style_header(timeline, 1, len(timeline_headers))
    for phase in TIMELINE_PHASES:
        timeline.append(list(phase))
        r = timeline.max_row
        for c in (2, 3):
            timeline.cell(row=r, column=c).number_format = "YYYY-MM-DD"
        fill = status_fill(str(phase[-1]))
        if fill:
            timeline.cell(row=r, column=7).fill = fill
    autosize(timeline, {1: 16, 2: 12, 3: 12, 4: 42, 5: 22, 6: 36, 7: 12})
    timeline.freeze_panes = "A2"

    # --- Branches sheet ---
    branches = wb.create_sheet("Git Branches")
    branch_headers = ["Branch Name", "Type", "Linked Story ID", "Title", "Target", "Status"]
    branches.append(branch_headers)
    style_header(branches, 1, len(branch_headers))
    for row in rows:
        story_id, _, _, _, _, title, _, branch, btype, _, status, *_ = row
        branches.append([branch, btype, story_id, title, "develop", status])
        fill = status_fill(status)
        if fill:
            branches.cell(row=branches.max_row, column=6).fill = fill
    autosize(branches, {1: 42, 2: 10, 3: 14, 4: 44, 5: 12, 6: 12})
    branches.freeze_panes = "A2"
    branches.auto_filter.ref = f"A1:F{branches.max_row}"

    # --- Workflow sheet (DoD) ---
    workflow = wb.create_sheet("Dev Workflow")
    workflow_headers = ["Step", "Stage", "Owner", "Input", "Output", "Tool / Gate"]
    workflow.append(workflow_headers)
    style_header(workflow, 1, len(workflow_headers))
    workflow_rows = [
        (1, "Story drafted", "PM", "Client requirement", "Story in Stories sheet", "PROJECT_TRACKING.xlsx"),
        (2, "Branch created", "Dev", "Story ID + title", "feat/fix branch", "Git"),
        (3, "Implementation", "Dev", "Acceptance criteria", "Code + tests", "IDE / Docker"),
        (4, "Migration (if schema)", "Backend", "Model change", "Alembic revision", "alembic revision"),
        (5, "Pull request", "Dev", "Branch commits", "PR to develop", "GitHub/GitLab"),
        (6, "Code review", "Peer", "PR", "Approved PR", "Review checklist"),
        (7, "CI / smoke", "CI", "PR build", "Green checks", "Docker compose tests"),
        (8, "Merge", "Maintainer", "Approved PR", "Updated develop", "Protected branch"),
        (9, "UAT on staging", "Client + QA", "Build on staging", "Pass/fail notes", "UAT_CHECKLIST_SHOPPING.md"),
        (10, "Admin↔Shopping verify", "QA", "WF-A to WF-F", "Reflection confirmed", "E2E Workflows sheet"),
        (11, "Story status update", "PM", "UAT result", "Status=Done", "Excel tracker"),
        (12, "Release to main", "DevOps", "Stable develop", "Production tag", "Release checklist"),
    ]
    for item in workflow_rows:
        workflow.append(list(item))
    autosize(workflow, {1: 8, 2: 24, 3: 14, 4: 22, 5: 24, 6: 28})

    # --- Integrations sheet ---
    integrations = wb.create_sheet("Integrations")
    int_headers = ["Integration", "Provider Options", "Phase", "Story IDs", "Key Objects", "Success Metric"]
    integrations.append(int_headers)
    style_header(integrations, 1, len(int_headers))
    integrations.append(
        [
            "Payments",
            "Razorpay (primary), Stripe (optional), COD",
            "P3 / S3",
            "VL-019, VL-020, VL-021, VS-017",
            "payments, payment_events, refunds",
            "Paid order confirmed via verified webhook",
        ]
    )
    integrations.append(
        [
            "Courier",
            "Shiprocket / Delhivery / BlueDart + manual",
            "P4 / S5",
            "VL-022, VL-023, VL-024, VS-024",
            "shipments, courier_accounts, labels",
            "AWB created from paid order < 2 minutes",
        ]
    )
    integrations.append(
        [
            "Tracking",
            "Partner webhooks + internal timeline",
            "P5 / S5",
            "VL-025, VL-026, VL-027, VL-028, VS-024, VS-025",
            "shipment_events, order_status_history, notifications",
            "Customer sees live timeline with AWB",
        ]
    )
    integrations.append(
        [
            "Catalog Sync",
            "Shared PostgreSQL products/categories/media",
            "S2",
            "VS-004, VS-005, VS-006, VS-007, VS-008, VS-009",
            "products, product_media, product_variants, brands",
            "Admin publish visible on shopping without redeploy",
        ]
    )
    autosize(integrations, {1: 14, 2: 48, 3: 12, 4: 42, 5: 48, 6: 48})

    # --- Dashboard sheet ---
    dash = wb.create_sheet("Dashboard", 0)
    dash["A1"] = "Valaiyagam E-Commerce — Project Tracking Dashboard"
    dash["A1"].font = Font(size=16, bold=True, color="0F766E")
    dash.merge_cells("A1:F1")

    dash["A3"] = "Project"
    dash["B3"] = "Valaiyagam Admin + Fashion Shopping (full sync)"
    dash["A4"] = "Start Date"
    dash["B4"] = START
    dash["B4"].number_format = "YYYY-MM-DD"
    dash["A5"] = "Planned End"
    dash["B5"] = week_date(16, 6)
    dash["B5"].number_format = "YYYY-MM-DD"
    dash["A6"] = "Document Date"
    dash["B6"] = date(2026, 7, 25)
    dash["B6"].number_format = "YYYY-MM-DD"
    dash["A7"] = "Architecture"
    dash["B7"] = "Admin :3000/:8000 · Shopping :3001/:8001 · PostgreSQL shared DB"

    done = sum(1 for r in rows if r[10] == "Done")
    in_progress = sum(1 for r in rows if r[10] == "In Progress")
    planned = sum(1 for r in rows if r[10] == "Planned")
    total = len(rows)
    vl_count = sum(1 for r in rows if str(r[0]).startswith("VL-"))
    vs_count = sum(1 for r in rows if str(r[0]).startswith("VS-"))

    dash["A9"] = "KPI"
    dash["B9"] = "Value"
    style_header(dash, 9, 2)
    dash["A10"] = "Total Stories (VL + VS)"
    dash["B10"] = total
    dash["A11"] = "Foundation stories (VL-*)"
    dash["B11"] = vl_count
    dash["A12"] = "Shopping sync stories (VS-*)"
    dash["B12"] = vs_count
    dash["A13"] = "Done"
    dash["B13"] = done
    dash["A14"] = "In Progress"
    dash["B14"] = in_progress
    dash["A15"] = "Planned"
    dash["B15"] = planned
    dash["A16"] = "% Complete"
    dash["B16"] = done / total if total else 0
    dash["B16"].number_format = "0.0%"

    dash["A18"] = "How to use this workbook"
    dash["A18"].font = Font(bold=True, size=12)
    instructions = [
        "1. Update Stories sheet status weekly (Done / In Progress / Planned / Blocked).",
        "2. Use E2E Workflows sheet for admin↔shopping reflection scenarios (WF-A to WF-F).",
        "3. Create Git branch using the exact name in column Git Feature / Branch.",
        "4. Link PRs to Story ID (example: VS-018) in the PR title.",
        "5. Review Timeline sheet in client weekly calls (P* + S* phases).",
        "6. Use Dev Workflow sheet as Definition of Done checklist.",
        "7. Run docs/UAT_CHECKLIST_SHOPPING.md for publish→buy→fulfill sign-off.",
        "8. Plan details: docs/SHOPPING_ADMIN_WORKFLOW_PLAN.md",
    ]
    for i, text in enumerate(instructions):
        dash.cell(row=19 + i, column=1, value=text)

    dash["A28"] = "Related docs"
    dash["A28"].font = Font(bold=True)
    dash["A29"] = "docs/SHOPPING_ADMIN_WORKFLOW_PLAN.md"
    dash["A30"] = "docs/UAT_CHECKLIST_SHOPPING.md"
    dash["A31"] = "docs/ECOM_ARCHITECTURE_AND_TIMELINE.md"
    dash["A32"] = "docs/architecture/system-architecture.md"
    dash["A33"] = "docs/api/api-documentation.md"

    autosize(dash, {1: 58, 2: 62})

    wb.save(path)


def main() -> None:
    pptx_path = ROOT / "Valaiyagam_Client_Presentation.pptx"
    xlsx_path = ROOT / "PROJECT_TRACKING.xlsx"
    build_pptx(pptx_path)
    build_excel(xlsx_path)
    print(f"Wrote {pptx_path}")
    print(f"Wrote {xlsx_path}")


if __name__ == "__main__":
    main()
